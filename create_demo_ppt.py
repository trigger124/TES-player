from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(10), Inches(1.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.85), Inches(9), Inches(1.1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(16)
        p.level = 0

    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Left column
    lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.4), Inches(0.6))
    tf = lt_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    lb_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.4), Inches(5))
    tf = lb_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)

    # Right column
    rt_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.4), Inches(0.6))
    tf = rt_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    rb_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.1), Inches(4.4), Inches(5))
    tf = rb_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "TES 比赛实时 AI 监听系统", "自动识别解说语义，第一时间感知胜负")

    add_section_slide(prs, "01 项目背景")

    add_content_slide(prs, "为什么要做这个系统？", [
        "看 TES 比赛时，希望第一时间知道比赛结果",
        "传统方案：人工盯屏、手动刷新网页，效率低",
        "现有工具：基于关键词匹配，容易漏判、误判",
        "需求：实时监听解说音频，用 AI 理解语义，自动判断比赛胜负",
        "目标：TES 赢了安静陪伴，输了立刻打开安慰视频"
    ])

    add_section_slide(prs, "02 系统架构")

    add_content_slide(prs, "数据流：从声音到决策", [
        "① 音频采集：pyaudiowpatch 捕获系统扬声器回环音频",
        "② VAD 人声检测：Silero VAD 实时判断是否在说话",
        "③ 语音识别：Whisper 把解说语音转成文字",
        "④ 语义理解：DeepSeek AI 分析解说内容，推断比赛状态",
        "⑤ 状态机：跟踪 initial → in_game → advantage/disadvantage → ending → win/lose",
        "⑥ 执行动作：胜利/失败时触发对应反馈"
    ])

    add_content_slide(prs, "核心技术栈", [
        "Python 3.12 + PyTorch 2.12.1 (CUDA 12.6)",
        "pyaudiowpatch：Windows 扬声器回环采集",
        "Silero VAD：轻量高效的人声检测",
        "OpenAI Whisper：本地语音识别（medium 模型）",
        "DeepSeek API：语义推理与比赛状态判断",
        "RTX 4060 Laptop GPU 加速推理"
    ])

    add_section_slide(prs, "03 核心亮点")

    add_two_column_slide(prs, "与传统方案对比", "传统关键词检测", [
        "依赖固定关键词库",
        "无法处理同义词、口语化表达",
        "容易漏听或误判",
        "只能判断单句，缺乏上下文"
    ], "本系统 AI 语义理解", [
        "DeepSeek 理解自然语言",
        "支持同义、反讽、口语表达",
        "维护比赛状态机，持续跟踪",
        "低置信度自动触发外部裁判复核"
    ])

    add_content_slide(prs, "实时性优化", [
        "录音 chunk 1024 帧，单帧延迟约 21ms",
        "VAD 缓冲到 512 样本再推理，避免 Input too short 错误",
        "torchaudio 专业重采样，避免混叠失真",
        "Whisper medium 模型跑在 RTX 4060 CUDA 上",
        "VAD 中段跳过检测，降低 CPU 占用"
    ])

    add_section_slide(prs, "04 演示流程")

    add_content_slide(prs, "Demo 步骤", [
        "1. 启动程序，加载 Whisper medium 模型",
        "2. 自动识别扬声器回环设备并开始录音",
        "3. 播放 TES 比赛解说/直播",
        "4. 实时看到 VAD 检测到语音、Whisper 识别文字",
        "5. DeepSeek 持续更新比赛状态与理由",
        "6. 当判断为失败时，自动打开一次安慰视频",
        "7. 当判断为胜利时，自动结束监听"
    ])

    add_section_slide(prs, "05 踩坑与解决")

    add_content_slide(prs, "开发中遇到的典型问题", [
        "Anaconda Python 与 torch 2.12.1 VC++ 运行时版本不匹配 → 换独立 Python 3.12",
        "Silero VAD 报 Input audio chunk is too short → 缓冲到 512 样本",
        "Whisper 默认 CPU 推理慢 → 启用 CUDA auto 设备",
        "检测到失败后连续打开多个网页 → 主线程统一加 game_ended 互斥标志",
        "重采样直接用抽点导致识别率低 → 改用 torchaudio Resample"
    ])

    add_section_slide(prs, "06 未来展望")

    add_content_slide(prs, "还可以做什么？", [
        "接入多路直播流，支持同时监听多场比赛",
        "加入弹幕/评论文本，多模态判断情绪与赛况",
        "做成桌面 GUI 或浏览器插件，降低使用门槛",
        "增加赛后数据回放与统计报表",
        "支持其他战队与更多游戏项目"
    ])

    add_title_slide(prs, "感谢观看", "项目地址 / 代码将在视频简介中提供")

    prs.save("TES_AI_Demo.pptx")
    print("PPT 已生成: TES_AI_Demo.pptx")


if __name__ == "__main__":
    main()
